"""
Micro-service PPTX/PDF — La Machine à CV
========================================

Reçoit le JSON envoyé par Lovable, ouvre le template `templates/{templateId}.pptx`,
remplace les zones de texte/images par les données utilisateur, exporte en PDF
via LibreOffice et renvoie deux URLs téléchargeables.

POST /generate
Auth: Bearer <PPTX_API_KEY>
Body: payload JSON Lovable (voir src/lib/pdf/payload.ts)
Resp: { pptxUrl, pdfUrl }

Démarrage local :
  uvicorn main:app --host 0.0.0.0 --port 8080

Déploiement Render/Railway/Fly : voir README.md
"""
from __future__ import annotations

import base64
import io
import os
import subprocess
import uuid
from pathlib import Path
from typing import Any

from fastapi import FastAPI, Header, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Emu
from pydantic import BaseModel, Field

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
ROOT = Path(__file__).parent
TEMPLATES_DIR = ROOT / "templates"
# Auto-création du dossier templates/ + migration depuis la racine si besoin
TEMPLATES_DIR.mkdir(exist_ok=True)
for _pptx in ROOT.glob("*.pptx"):
    target = TEMPLATES_DIR / _pptx.name
    if not target.exists():
        try:
            _pptx.replace(target)
        except Exception:
            # fallback: copie
            import shutil
            shutil.copy2(_pptx, target)
# Emplacements alternatifs supportés
TEMPLATE_SEARCH_DIRS = [TEMPLATES_DIR, ROOT]
OUTPUT_DIR = ROOT / "out"
OUTPUT_DIR.mkdir(exist_ok=True)


def _find_template(template_id: str) -> Path | None:
    name = f"{template_id}.pptx"
    for d in TEMPLATE_SEARCH_DIRS:
        p = d / name
        if p.exists():
            return p
    return None


def _list_templates() -> list[str]:
    seen: set[str] = set()
    out: list[str] = []
    for d in TEMPLATE_SEARCH_DIRS:
        if not d.exists():
            continue
        for p in sorted(d.glob("*.pptx")):
            if p.stem not in seen:
                seen.add(p.stem)
                out.append(p.stem)
    return out

API_KEY = os.environ.get("PPTX_API_KEY", "")
PUBLIC_BASE_URL = os.environ.get("PUBLIC_BASE_URL", "http://localhost:8080")

app = FastAPI(title="Machine à CV — PPTX service")
app.add_middleware(
    CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"]
)

# ---------------------------------------------------------------------------
# Payload
# ---------------------------------------------------------------------------
class Experience(BaseModel):
    id: str = ""
    title: str = ""
    company: str = ""
    period: str = ""
    description: str = ""


class Education(BaseModel):
    id: str = ""
    title: str = ""
    school: str = ""
    period: str = ""


class SoftSkill(BaseModel):
    label: str = ""
    def_: str = Field("", alias="def")

    class Config:
        populate_by_name = True


class Language(BaseModel):
    label: str = ""
    level: str = ""


class CvPayload(BaseModel):
    templateId: str = "cv-1"
    fullName: str = ""
    firstName: str = ""
    lastName: str = ""
    jobTitle: str = ""
    email: str = ""
    phone: str = ""
    city: str = ""
    address: str = ""
    summary: str = ""
    photo: str | None = None  # data URL ou base64
    experiences: list[Experience] = []
    education: list[Education] = []
    skills: list[str] = []
    softSkills: list[SoftSkill] = []
    languages: list[Language] = []
    interests: list[str] = []
    offerAnalysis: Any = None
    atsScore: Any = None
    createdAt: str = ""


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _decode_photo(photo: str | None) -> bytes | None:
    if not photo:
        return None
    if photo.startswith("data:"):
        try:
            photo = photo.split(",", 1)[1]
        except IndexError:
            return None
    try:
        return base64.b64decode(photo)
    except Exception:
        return None


def _fmt_experiences(items: list[Experience]) -> str:
    blocks = []
    for x in items:
        head = " — ".join(filter(None, [x.title, x.company]))
        line = f"{head}\n{x.period}\n{x.description}".strip()
        blocks.append(line)
    return "\n\n".join(blocks)


def _fmt_education(items: list[Education]) -> str:
    return "\n".join(
        " — ".join(filter(None, [x.title, x.school, x.period])) for x in items
    )


def _build_replacement_map(p: CvPayload) -> dict[str, str]:
    """
    Mappe les jetons {{TOKEN}} présents dans le PPTX vers la donnée user.
    Édite ton template CV 1.pptx pour utiliser ces jetons une fois, puis ce
    service les remplacera automatiquement à chaque génération.
    """
    return {
        "{{FULL_NAME}}": p.fullName or f"{p.firstName} {p.lastName}".strip(),
        "{{FIRST_NAME}}": p.firstName,
        "{{LAST_NAME}}": p.lastName,
        "{{JOB_TITLE}}": p.jobTitle,
        "{{EMAIL}}": p.email,
        "{{PHONE}}": p.phone,
        "{{CITY}}": p.city,
        "{{ADDRESS}}": p.address,
        "{{SUMMARY}}": p.summary,
        "{{EXPERIENCES}}": _fmt_experiences(p.experiences),
        "{{EDUCATION}}": _fmt_education(p.education),
        "{{SKILLS}}": " · ".join(p.skills),
        "{{SOFT_SKILLS}}": " · ".join(s.label for s in p.softSkills),
        "{{LANGUAGES}}": " · ".join(
            " ".join(filter(None, [l.label, f"({l.level})" if l.level else ""])) for l in p.languages
        ),
        "{{INTERESTS}}": " · ".join(p.interests),
    }


def _replace_text_in_runs(shape, mapping: dict[str, str]) -> None:
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for token, value in mapping.items():
                if token in run.text:
                    run.text = run.text.replace(token, value or "")


def _replace_image_placeholder(shape, photo_bytes: bytes) -> bool:
    """Si la shape contient {{PHOTO}}, on la remplace par l'image."""
    if not shape.has_text_frame or "{{PHOTO}}" not in shape.text_frame.text:
        return False
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    slide = shape._parent  # type: ignore[attr-defined]
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(io.BytesIO(photo_bytes), left, top, width=width, height=height)
    return True


def render_pptx(payload: CvPayload) -> Path:
    template_file = _find_template(payload.templateId)
    if template_file is None:
        available = _list_templates()
        raise HTTPException(
            404,
            f"Template introuvable: {payload.templateId}. Disponibles: {available}",
        )

    prs = Presentation(template_file)
    mapping = _build_replacement_map(payload)
    photo_bytes = _decode_photo(payload.photo)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if photo_bytes and _replace_image_placeholder(shape, photo_bytes):
                continue
            _replace_text_in_runs(shape, mapping)
            # tableaux
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for token, value in mapping.items():
                                    if token in run.text:
                                        run.text = run.text.replace(token, value or "")

    out_id = uuid.uuid4().hex
    out_pptx = OUTPUT_DIR / f"{out_id}.pptx"
    prs.save(out_pptx)
    return out_pptx


def convert_to_pdf(pptx_path: Path) -> Path:
    """Conversion via LibreOffice (`soffice` doit être installé)."""
    cmd = [
        "soffice", "--headless", "--convert-to", "pdf",
        "--outdir", str(OUTPUT_DIR), str(pptx_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=90)
    if result.returncode != 0:
        raise HTTPException(500, f"LibreOffice failed: {result.stderr}")
    pdf_path = pptx_path.with_suffix(".pdf")
    if not pdf_path.exists():
        raise HTTPException(500, "PDF non généré")
    return pdf_path


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------
@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.get("/templates")
def list_templates() -> dict[str, Any]:
    return {
        "templates": _list_templates(),
        "searchDirs": [str(d) for d in TEMPLATE_SEARCH_DIRS],
    }


@app.post("/generate")
def generate(payload: CvPayload, authorization: str | None = Header(default=None)):
    if API_KEY:
        if not authorization or not authorization.startswith("Bearer "):
            raise HTTPException(401, "Missing bearer")
        if authorization.split(" ", 1)[1] != API_KEY:
            raise HTTPException(401, "Invalid key")

    pptx_path = render_pptx(payload)
    pdf_path = convert_to_pdf(pptx_path)
    return {
        "pptxUrl": f"{PUBLIC_BASE_URL}/files/{pptx_path.name}",
        "pdfUrl": f"{PUBLIC_BASE_URL}/files/{pdf_path.name}",
    }


@app.get("/files/{name}")
def download(name: str):
    f = OUTPUT_DIR / name
    if not f.exists() or ".." in name:
        raise HTTPException(404, "Not found")
    return FileResponse(f, filename=name)
